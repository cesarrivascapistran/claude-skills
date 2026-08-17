"""
Deck Builder - build_deck.py

Builds PowerPoint presentations from a JSON specification, on top of whatever
.pptx template you point it at. No brand assets are bundled: colors come from a
tokens file and the template is supplied by the caller.

Usage:
    python build_deck.py --spec deck_spec.json --template /path/to/template.pptx

Template resolution order:
    1. --template CLI argument
    2. DECK_TEMPLATE_PATH environment variable
    3. CEMEX_TEMPLATE_PATH environment variable (kept for existing setups)
    4. template.pptx next to the spec file
"""

import argparse
import json
import os
import re
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── Month names for language-aware date fallback ───────────────────────────────
MONTHS_ES = [
    "", "Enero", "Febrero", "Marzo", "Abril", "Mayo", "Junio",
    "Julio", "Agosto", "Septiembre", "Octubre", "Noviembre", "Diciembre",
]

# ── Template path resolution ──────────────────────────────────────────────────
_SCRIPT_DIR = Path(__file__).resolve().parent          # skills/deck/scripts/
_SKILL_DIR = _SCRIPT_DIR.parent                        # skills/deck/
_WORKSPACE_ROOT = _SKILL_DIR.parent.parent.parent      # workspace root

def resolve_template_path(cli_path: str = None) -> str:
    """Resolve the .pptx template using the documented fallback chain."""
    # 1. Explicit CLI argument
    if cli_path and os.path.exists(cli_path):
        return cli_path
    # 2. Environment variable
    env_path = os.environ.get("DECK_TEMPLATE_PATH") or os.environ.get("CEMEX_TEMPLATE_PATH")
    if env_path and os.path.exists(env_path):
        return env_path
    # 3. Workspace root convention
    workspace_path = _WORKSPACE_ROOT / "template.pptx"
    if workspace_path.exists():
        return str(workspace_path)
    # 4. Return the workspace path as default (will fail at build time with clear error)
    return str(workspace_path)

# ── Design Tokens ────────────────────────────────────────────────────────
# Neutral defaults on purpose: no organization's palette is baked in. Override
# every value with --tokens tokens.json (see references/brand_tokens.example.json).
# Most slide chrome comes from the template masters, so these mainly drive chart
# series and emphasis text.
COLORS = {
    "dark_text":      RGBColor(0x3D, 0x3D, 0x3D),
    "white":          RGBColor(0xFF, 0xFF, 0xFF),
    "brand_primary":  RGBColor(0x1F, 0x4E, 0x79),
    "brand_accent":   RGBColor(0x2E, 0x8B, 0x57),
    "blue_light":     RGBColor(0x4A, 0x90, 0xD9),
    "red":            RGBColor(0xC0, 0x39, 0x2B),
    "green_soft":     RGBColor(0x5C, 0xB8, 0x5C),
    "amber":          RGBColor(0xE6, 0xA2, 0x3A),
    "orange":         RGBColor(0xE0, 0x7B, 0x39),
    "purple":         RGBColor(0x7B, 0x5E, 0xA7),
    "brand_deep":     RGBColor(0x16, 0x36, 0x54),
}


def load_color_tokens(tokens_path: str) -> None:
    """Override COLORS in place from a tokens JSON file.

    Accepts either {"colors": {"key": "#RRGGBB"}} or
    {"colors": {"key": {"hex": "#RRGGBB"}}}. Unknown keys are ignored so a
    tokens file can carry extra metadata without breaking the build.
    """
    with open(tokens_path, "r", encoding="utf-8") as fh:
        data = json.load(fh)
    colors = data.get("colors", data)
    applied = 0
    for key, value in colors.items():
        if key not in COLORS:
            continue
        hex_str = value.get("hex") if isinstance(value, dict) else value
        if not isinstance(hex_str, str):
            continue
        COLORS[key] = hex_to_rgb(hex_str)
        applied += 1
    # CHART_COLORS captured the originals by value, so rebuild it in place.
    CHART_COLORS[:] = [
        COLORS["brand_primary"],
        COLORS["brand_accent"],
        COLORS["blue_light"],
        COLORS["amber"],
        COLORS["red"],
        COLORS["orange"],
        COLORS["purple"],
        COLORS["green_soft"],
    ]
    print(f"Loaded {applied} color tokens from {tokens_path}")

CHART_COLORS = [
    COLORS["brand_primary"],
    COLORS["brand_accent"],
    COLORS["blue_light"],
    COLORS["amber"],
    COLORS["red"],
    COLORS["orange"],
    COLORS["purple"],
    COLORS["green_soft"],
]

FONT_NAME = "Arial"

# ── Layout resolution ──────────────────────────────────────────────────────────
# Logical key -> ordered list of candidate layout names to match in the template.
# The first match wins. This makes the builder resilient to templates that ship
# with a subset of layouts (e.g. only "Cover 4" but no "Cover 1").
LAYOUT_NAME_CANDIDATES = {
    "cover_1":        ["Cover 1", "Cover 4"],
    "cover_4":        ["Cover 4", "Cover 1"],
    "agenda":         ["Agenda"],
    "section_header": ["Section Header"],
    "title_only":     ["Title Only"],
    "title_and_text": ["Title and Text"],
    "two_column":     ["Two Text Column & Title", "Two Column", "Two Text Column"],
    "content":        ["Content Layout", "Content"],
    "blank":          ["Blank"],
    "thank_you":      ["Thank You"],
    "backup":         ["Back Up Slides - Blue", "Back Up Slides", "Backup", "Blank"],
}

# Populated by `resolve_layouts(prs)` at build time.
LAYOUT: dict[str, int] = {}


def resolve_layouts(prs: Presentation) -> dict:
    """Resolve logical layout keys to real indices using layout names.

    Falls back through candidate names so decks render even when the template
    has a reduced layout set. Raises if a required layout cannot be found.
    """
    name_to_idx = {layout.name: i for i, layout in enumerate(prs.slide_layouts)}
    resolved: dict[str, int] = {}
    missing: list[str] = []

    for key, candidates in LAYOUT_NAME_CANDIDATES.items():
        for candidate in candidates:
            if candidate in name_to_idx:
                resolved[key] = name_to_idx[candidate]
                break
        else:
            missing.append(key)

    if missing:
        available = ", ".join(f"{i}:{n}" for n, i in name_to_idx.items())
        raise RuntimeError(
            f"Template is missing required layouts for keys: {missing}. "
            f"Available layouts: {available}"
        )

    return resolved

CHART_TYPE_MAP = {
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_clustered":    XL_CHART_TYPE.BAR_CLUSTERED,
    "line":             XL_CHART_TYPE.LINE,
    "pie":              XL_CHART_TYPE.PIE,
    "doughnut":         XL_CHART_TYPE.DOUGHNUT,
    "area":             XL_CHART_TYPE.AREA,
}


# ── Helpers ────────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_str: str) -> RGBColor:
    """Convert a hex color string (with or without #) to RGBColor."""
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def set_text(placeholder, text: str, font_size: int = None, bold: bool = None,
             color: RGBColor = None, alignment=None):
    """Set text in a placeholder with optional formatting."""
    placeholder.text = text
    for paragraph in placeholder.text_frame.paragraphs:
        if alignment:
            paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            if font_size:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color


def add_bullets(text_frame, items: list, font_size: int = 18,
                color: RGBColor = None, bold: bool = False):
    """Add bullet points to a text frame. First item replaces existing, rest appended."""
    text_color = color or COLORS["dark_text"]

    for i, item in enumerate(items):
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()

        para.text = item
        para.level = 0
        para.space_before = Pt(6)
        para.space_after = Pt(2)

        for run in para.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = text_color


def add_table_to_slide(slide, table_spec: dict, position: dict = None):
    """Add a formatted table to a slide."""
    data = table_spec["data"]
    rows = len(data)
    cols = len(data[0]) if rows > 0 else 0

    col_widths = table_spec.get("col_widths")
    pos = position or {}
    x = Inches(pos.get("x", 0.5))
    y = Inches(pos.get("y", 1.5))

    # Auto-calculate width from col_widths sum if provided
    if col_widths:
        calculated_width = sum(col_widths[:cols])
        width = Inches(pos.get("width", calculated_width))
    else:
        width = Inches(pos.get("width", 12.0))
    height = Inches(pos.get("height", 4.5))

    table_shape = slide.shapes.add_table(rows, cols, x, y, width, height)
    table = table_shape.table

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            if i < cols:
                table.columns[i].width = Inches(w)
        # Warn if col_widths don't match table width
        cw_sum = sum(col_widths[:cols])
        table_w = round(width / 914400, 2)
        if abs(cw_sum - table_w) > 0.5:
            print(f"  WARNING: col_widths sum ({cw_sum}) differs from table width ({table_w})")

    header_bg = hex_to_rgb(table_spec.get("header_bg", "0000B3"))
    header_font_color = hex_to_rgb(table_spec.get("header_font_color", "FFFFFF"))
    alt_row_bg = table_spec.get("alt_row_bg")

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Format paragraph
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(11)

            # Header row styling
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.bold = True
                        run.font.color.rgb = header_font_color
                        run.font.size = Pt(12)
            elif alt_row_bg and row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(alt_row_bg)

    return table_shape


def add_chart_to_slide(slide, chart_spec: dict):
    """Add a chart to a slide based on specification."""
    chart_type_key = chart_spec.get("type", "column_clustered")
    chart_type = CHART_TYPE_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = chart_spec["categories"]

    for series in chart_spec.get("series", []):
        chart_data.add_series(series["name"], series["values"])

    pos = chart_spec.get("position", {})
    x = Inches(pos.get("x", 0.5))
    y = Inches(pos.get("y", 1.5))
    cx = Inches(pos.get("width", 12.0))
    cy = Inches(pos.get("height", 5.5))

    graphic_frame = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = len(chart_spec.get("series", [])) > 1

    # Apply palette colors to series
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = CHART_COLORS[i % len(CHART_COLORS)]

    return graphic_frame


def add_images_to_slide(slide, images: list):
    """Add images to a slide from a list of image specs."""
    for img in images:
        path = img.get("path", "")
        if not os.path.exists(path):
            print(f"WARNING: Image not found: {path}")
            continue
        x = Inches(img.get("x", 0.5))
        y = Inches(img.get("y", 1.5))
        width = Inches(img.get("width")) if "width" in img else None
        height = Inches(img.get("height")) if "height" in img else None
        slide.shapes.add_picture(path, x, y, width, height)


# ── Slide Builders ─────────────────────────────────────────────────────────────

def build_cover(prs: Presentation, spec: dict, language: str = "es"):
    """Build a cover slide."""
    layout_idx = spec.get("layout", LAYOUT["cover_1"])
    if isinstance(layout_idx, str):
        layout_idx = LAYOUT.get(layout_idx, LAYOUT["cover_1"])

    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    # idx 0: Title
    if 0 in {ph.placeholder_format.idx for ph in slide.placeholders}:
        set_text(slide.placeholders[0], spec.get("title", ""),
                 font_size=36, bold=True, color=COLORS["brand_primary"])

    # idx 12: Description
    if 12 in {ph.placeholder_format.idx for ph in slide.placeholders}:
        set_text(slide.placeholders[12], spec.get("description", ""),
                 font_size=18, color=COLORS["dark_text"])

    # idx 15: Presenter
    if 15 in {ph.placeholder_format.idx for ph in slide.placeholders}:
        set_text(slide.placeholders[15], spec.get("presenter", ""),
                 font_size=12, color=COLORS["dark_text"])

    # idx 14: Date (language-aware fallback)
    if 14 in {ph.placeholder_format.idx for ph in slide.placeholders}:
        date_text = spec.get("date")
        if not date_text:
            now = datetime.now()
            if language == "es":
                date_text = f"{MONTHS_ES[now.month]} {now.year}"
            else:
                date_text = now.strftime("%B %Y")
        set_text(slide.placeholders[14], date_text,
                 font_size=12, color=COLORS["dark_text"])

    return slide


def build_agenda(prs: Presentation, items: list):
    """Build an agenda slide with numbered items."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT["agenda"]])

    # idx 12: Agenda body
    ph_indices = {ph.placeholder_format.idx for ph in slide.placeholders}
    if 12 in ph_indices:
        tf = slide.placeholders[12].text_frame
        for i, item in enumerate(items):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            para.text = f"{i + 1:02d}   {item}"
            para.space_before = Pt(12)
            para.space_after = Pt(6)
            para.level = 0

            for run in para.runs:
                run.font.name = FONT_NAME
                run.font.size = Pt(20)
                run.font.color.rgb = COLORS["brand_deep"]
                run.font.bold = False

    return slide


def build_section_header(prs: Presentation, spec: dict):
    """Build a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT["section_header"]])

    ph_indices = {ph.placeholder_format.idx for ph in slide.placeholders}

    # idx 10: Section title (e.g., "01")
    if 10 in ph_indices:
        set_text(slide.placeholders[10], spec.get("title", ""),
                 font_size=32, bold=True, color=COLORS["brand_primary"])

    # idx 11: Section subtitle
    if 11 in ph_indices:
        set_text(slide.placeholders[11], spec.get("subtitle", ""),
                 font_size=20, color=COLORS["dark_text"])

    return slide


def build_title_and_text(prs: Presentation, spec: dict):
    """Build a title + subtitle + body text slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT["title_and_text"]])

    ph_indices = {ph.placeholder_format.idx for ph in slide.placeholders}

    # idx 0: Title
    if 0 in ph_indices:
        set_text(slide.placeholders[0], spec.get("title", ""),
                 font_size=28, bold=True, color=COLORS["brand_primary"])

    # idx 10: Subtitle
    if 10 in ph_indices:
        set_text(slide.placeholders[10], spec.get("subtitle", ""),
                 font_size=14, color=COLORS["dark_text"])

    # idx 14: Body
    if 14 in ph_indices:
        body = spec.get("body", [])
        if isinstance(body, list):
            add_bullets(slide.placeholders[14].text_frame, body,
                       font_size=18, color=COLORS["dark_text"])
        else:
            set_text(slide.placeholders[14], str(body),
                     font_size=18, color=COLORS["dark_text"])

    # Optional table
    if "table" in spec:
        table_pos = {"x": 0.5, "y": 2.5, "width": 12.0, "height": 4.0}
        add_table_to_slide(slide, spec["table"], table_pos)

    # Optional chart
    if "chart" in spec:
        add_chart_to_slide(slide, spec["chart"])

    return slide


def build_two_column(prs: Presentation, spec: dict):
    """Build a two-column slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT["two_column"]])

    ph_indices = {ph.placeholder_format.idx for ph in slide.placeholders}

    # idx 0: Title
    if 0 in ph_indices:
        set_text(slide.placeholders[0], spec.get("title", ""),
                 font_size=28, bold=True, color=COLORS["brand_primary"])

    # idx 10: Subtitle
    if 10 in ph_indices:
        set_text(slide.placeholders[10], spec.get("subtitle", ""),
                 font_size=14, color=COLORS["dark_text"])

    # idx 19: Left column title
    if 19 in ph_indices:
        set_text(slide.placeholders[19], spec.get("left_title", ""),
                 font_size=16, bold=True, color=COLORS["brand_primary"])

    # idx 20: Right column title
    if 20 in ph_indices:
        set_text(slide.placeholders[20], spec.get("right_title", ""),
                 font_size=16, bold=True, color=COLORS["brand_primary"])

    # idx 17: Left body
    if 17 in ph_indices:
        left_body = spec.get("left_body", [])
        if isinstance(left_body, list):
            add_bullets(slide.placeholders[17].text_frame, left_body,
                       font_size=16, color=COLORS["dark_text"])
        else:
            set_text(slide.placeholders[17], str(left_body),
                     font_size=16, color=COLORS["dark_text"])

    # idx 21: Right body
    if 21 in ph_indices:
        right_body = spec.get("right_body", [])
        if isinstance(right_body, list):
            add_bullets(slide.placeholders[21].text_frame, right_body,
                       font_size=16, color=COLORS["dark_text"])
        else:
            set_text(slide.placeholders[21], str(right_body),
                     font_size=16, color=COLORS["dark_text"])

    return slide


def build_kpi_callout(prs: Presentation, spec: dict):
    """Build a KPI callout slide with large centered metric numbers."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT["title_only"]])

    ph_indices = {ph.placeholder_format.idx for ph in slide.placeholders}

    if 0 in ph_indices:
        set_text(slide.placeholders[0], spec.get("title", ""),
                 font_size=28, bold=True, color=COLORS["brand_primary"])

    if 10 in ph_indices:
        set_text(slide.placeholders[10], spec.get("subtitle", ""),
                 font_size=14, color=COLORS["dark_text"])

    metrics = spec.get("metrics", [])
    count = len(metrics)
    if count == 0:
        return slide

    # Distribute metric cards evenly across the slide width
    slide_width = 13.33
    margin = 0.5
    usable_width = slide_width - 2 * margin
    card_width = min(usable_width / count, 4.0)
    total_cards_width = card_width * count
    start_x = (slide_width - total_cards_width) / 2

    card_top = 2.2
    card_height = 4.0

    for i, metric in enumerate(metrics):
        x = start_x + i * card_width
        color_name = metric.get("color", "brand_primary")
        color = COLORS.get(color_name, COLORS["brand_primary"])

        # Big number
        value_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(card_top),
            Inches(card_width - 0.4), Inches(card_height * 0.6)
        )
        vtf = value_box.text_frame
        vtf.word_wrap = True
        vp = vtf.paragraphs[0]
        vp.text = str(metric.get("value", ""))
        vp.alignment = PP_ALIGN.CENTER
        vp.space_before = Pt(50)
        for run in vp.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(60)
            run.font.bold = True
            run.font.color.rgb = color

        # Label below the number
        label_box = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(card_top + card_height * 0.6),
            Inches(card_width - 0.4), Inches(card_height * 0.35)
        )
        ltf = label_box.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.text = str(metric.get("label", ""))
        lp.alignment = PP_ALIGN.CENTER
        for run in lp.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(16)
            run.font.bold = False
            run.font.color.rgb = COLORS["dark_text"]

    return slide


def build_title_only(prs: Presentation, spec: dict):
    """Build a title-only slide (canvas for tables, charts, custom shapes)."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT["title_only"]])

    ph_indices = {ph.placeholder_format.idx for ph in slide.placeholders}

    if 0 in ph_indices:
        set_text(slide.placeholders[0], spec.get("title", ""),
                 font_size=28, bold=True, color=COLORS["brand_primary"])

    if 10 in ph_indices:
        set_text(slide.placeholders[10], spec.get("subtitle", ""),
                 font_size=14, color=COLORS["dark_text"])

    if "table" in spec:
        add_table_to_slide(slide, spec["table"])

    if "chart" in spec:
        add_chart_to_slide(slide, spec["chart"])

    return slide


def build_blank(prs: Presentation, spec: dict):
    """Build a blank slide with optional custom elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT["blank"]])

    if "table" in spec:
        add_table_to_slide(slide, spec["table"])

    if "chart" in spec:
        add_chart_to_slide(slide, spec["chart"])

    return slide


def build_thank_you(prs: Presentation):
    """Build the closing Thank You slide."""
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT["thank_you"]])


def build_backup_header(prs: Presentation):
    """Build the Back Up Slides divider."""
    return prs.slides.add_slide(prs.slide_layouts[LAYOUT["backup"]])


# ── Slide dispatcher ───────────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "title_and_text": build_title_and_text,
    "two_column":     build_two_column,
    "title_only":     build_title_only,
    "content":        build_title_only,
    "blank":          build_blank,
    "kpi_callout":    build_kpi_callout,
}


def build_content_slide(prs: Presentation, slide_spec: dict):
    """Dispatch to the appropriate slide builder based on layout type."""
    layout_type = slide_spec.get("layout", "title_and_text")
    builder = SLIDE_BUILDERS.get(layout_type, build_title_and_text)
    slide = builder(prs, slide_spec)

    # Add images if specified
    if "images" in slide_spec:
        add_images_to_slide(slide, slide_spec["images"])

    return slide


# ── Main orchestrator ──────────────────────────────────────────────────────────

def build_deck(spec: dict, output_path: str = None, template_path: str = None) -> str:
    """Build a complete deck from a JSON specification.

    Args:
        spec: Deck specification dictionary.
        output_path: Destination .pptx path (auto-generated if omitted).
        template_path: Explicit path to the .pptx template file.

    Returns the path to the generated PPTX file.
    """
    resolved_template = resolve_template_path(template_path)
    if not os.path.exists(resolved_template):
        print(f"ERROR: Template not found. Searched:\n"
              f"  1. --template CLI arg\n"
              f"  2. DECK_TEMPLATE_PATH or CEMEX_TEMPLATE_PATH env var\n"
              f"  3. {_WORKSPACE_ROOT / 'template.pptx'}\n"
              f"Set DECK_TEMPLATE_PATH or pass --template to fix.")
        sys.exit(1)

    prs = Presentation(resolved_template)

    # Resolve logical layout keys against this template's actual layout names.
    LAYOUT.clear()
    LAYOUT.update(resolve_layouts(prs))

    # Remove existing slides (template comes with sample slides)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
        )
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    language = spec.get("language", "es")

    # 1. Opening Cover (defaults to Cover 1)
    cover_spec = spec.get("cover", {})
    if cover_spec:
        if "layout" not in cover_spec:
            cover_spec["layout"] = LAYOUT["cover_1"]
        cover_slide = build_cover(prs, cover_spec, language)
        if "images" in cover_spec:
            add_images_to_slide(cover_slide, cover_spec["images"])

    # 2. Agenda
    agenda_items = spec.get("agenda", [])
    if agenda_items:
        build_agenda(prs, agenda_items)

    # 3. Sections
    for section in spec.get("sections", []):
        # Section header
        header = section.get("header", {})
        if header:
            build_section_header(prs, header)

        # Section content slides
        for slide_spec in section.get("slides", []):
            build_content_slide(prs, slide_spec)

    # 4. Closing Cover (defaults to Cover 4, before Thank You)
    closing_cover = spec.get("closing_cover")
    if closing_cover:
        if "layout" not in closing_cover:
            closing_cover["layout"] = LAYOUT["cover_4"]
        closing_slide = build_cover(prs, closing_cover, language)
        if "images" in closing_cover:
            add_images_to_slide(closing_slide, closing_cover["images"])

    # 5. Thank You
    if spec.get("thank_you", True):
        build_thank_you(prs)

    # 6. Backup slides
    if spec.get("backup", False):
        build_backup_header(prs)
        for slide_spec in spec.get("backup_slides", []):
            build_content_slide(prs, slide_spec)

    # Generate output path
    if not output_path:
        title = spec.get("cover", {}).get("title", "Presentation")
        safe_title = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')[:50]
        output_dir = Path("output")
        output_dir.mkdir(exist_ok=True)
        output_path = str(output_dir / f"{safe_title}.pptx")

    prs.save(output_path)
    slide_count = len(prs.slides)
    print(f"Presentation saved: {output_path} ({slide_count} slides)")
    return output_path


# ── CLI ────────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Deck Builder")
    parser.add_argument("--spec", required=True, help="Path to JSON specification file")
    parser.add_argument("--output", default=None, help="Output PPTX path")
    parser.add_argument("--template", default=None,
                        help="Path to the .pptx template (overrides env var and auto-detect)")
    parser.add_argument("--tokens", default=None,
                        help="Path to a colors JSON file overriding the neutral defaults")
    args = parser.parse_args()

    if args.tokens:
        load_color_tokens(args.tokens)

    with open(args.spec, "r", encoding="utf-8") as f:
        spec = json.load(f)

    build_deck(spec, args.output, args.template)


if __name__ == "__main__":
    main()
